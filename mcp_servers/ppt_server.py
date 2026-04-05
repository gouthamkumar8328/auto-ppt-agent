from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN
from pptx.enum.shapes import MSO_SHAPE, MSO_CONNECTOR
from pptx.dml.color import RGBColor
from datetime import datetime
from agent.diagram_agent import generate_diagram
import os
import requests
import glob


class PPTServer:
    """
    MCP Tool for creating and managing PowerPoint presentations.

    This class acts as a tool that allows an AI agent to:
    - Create a new presentation
    - Add slides with different layouts (text, visual, mixed, diagram)
    - Save the presentation to disk

    It encapsulates all PowerPoint-related operations.
    """

    def __init__(self):
        """
        Initialize a new PowerPoint presentation instance.

        Creates a blank presentation using python-pptx.
        """
        self.prs = Presentation()

    def add_slide(self, title, bullets, slide_type, topic):
        """
        Add a slide to the presentation based on the given type.

        Args:
            title (str): Title of the slide.
            bullets (list): List of bullet points (3–5 recommended).
            slide_type (str): Type of slide layout.
                Supported types:
                - "text": Text-focused slide
                - "visual": Image-focused slide
                - "mixed": Text + image
                - "diagram": Flow/process diagram
            topic (str): Presentation topic (used for images/diagrams).

        Behavior:
            - Creates a blank slide layout
            - Adds title
            - Adds content depending on slide type
        """
        slide = self.prs.slides.add_slide(self.prs.slide_layouts[6])

        # Add title
        self._add_title(slide, title)

        # Add content based on type
        if slide_type == "text":
            self._add_text_card(slide, bullets)
            self._add_image(slide, topic)   
            self._add_highlight(slide, bullets)

        elif slide_type == "visual":
            self._add_image(slide, topic)

        elif slide_type == "mixed":
            self._add_text_card(slide, bullets)
            self._add_image(slide, topic)

        elif slide_type == "diagram":
            steps = generate_diagram(topic, title)
            if steps:
                self._add_flow_diagram(slide, steps)
            else:
                self._add_text_card(slide, bullets)

    def _add_title(self, slide, title):
        """
        Add a styled title textbox to the slide.

        Args:
            slide: PowerPoint slide object
            title (str): Title text
        """
        box = slide.shapes.add_textbox(
            Inches(1), Inches(0.5), Inches(8), Inches(1)
        )

        tf = box.text_frame
        p = tf.paragraphs[0]

        p.text = title
        p.font.size = Pt(40)
        p.font.bold = True
        p.font.color.rgb = RGBColor(0, 90, 140)
        p.alignment = PP_ALIGN.LEFT

    def _add_text_card(self, slide, bullets):
        """
        Add a styled text card with bullet points.

        Args:
            slide: PowerPoint slide object
            bullets (list): List of bullet strings
        """
        card = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(1),
            Inches(2),
            Inches(6),
            Inches(3.5)
        )

        card.fill.solid()
        card.fill.fore_color.rgb = RGBColor(245, 248, 255)
        card.line.fill.background()

        tf = card.text_frame
        tf.clear()

        for i, b in enumerate(bullets[:4]):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()

            p.text = b
            p.font.size = Pt(20)
            p.space_after = Pt(10)
            p.font.bold = True
            p.font.color.rgb = RGBColor(30, 30, 30)

    def _add_side_panel(self, slide):
        """
        Add a decorative side panel to improve slide design.

        Args:
            slide: PowerPoint slide object
        """
        panel = slide.shapes.add_shape(
            MSO_SHAPE.RECTANGLE,
            Inches(7),
            Inches(1.5),
            Inches(2.5),
            Inches(4.5)
        )

        panel.fill.solid()
        panel.fill.fore_color.rgb = RGBColor(210, 225, 255)
        panel.line.fill.background()

    def _add_highlight(self, slide, bullets):
        """
        Add a highlighted key idea box.

        Args:
            slide: PowerPoint slide object
            bullets (list): Bullet points (used to extract key idea)
        """
        text = bullets[0] if bullets else "Key Idea"

        box = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE,
            Inches(7),
            Inches(4),
            Inches(2.5),
            Inches(1)
        )

        box.fill.solid()
        box.fill.fore_color.rgb = RGBColor(255, 235, 200)
        box.line.fill.background()

        tf = box.text_frame
        p = tf.paragraphs[0]

        p.text = f"⭐ {text}"
        p.font.size = Pt(20)

    def _add_image(self, slide, topic):
        """
        Add an image related to the topic using Unsplash.

        Args:
            slide: PowerPoint slide object
            topic (str): Search keyword for image
        """
        img = self._download_image(topic)

        if img:
            slide.shapes.add_picture(
                img,
                Inches(5.5),
                Inches(2),
                width=Inches(3.5)
            )

    def _add_flow_diagram(self, slide, steps):
        """
        Add a simple flow diagram using shapes and connectors.

        Args:
            slide: PowerPoint slide object
            steps (list): Steps in the process
        """
        steps = steps[:5]

        total_width = len(steps) * 1.6
        start_x = max(0.5, (10 - total_width) / 2)

        left = Inches(start_x)
        top = Inches(5)

        shapes = []

        for i, step in enumerate(steps):
            box = slide.shapes.add_shape(
                MSO_SHAPE.ROUNDED_RECTANGLE,
                left + Inches(i * 1.6),
                top,
                Inches(1.5),
                Inches(0.9)
            )

            box.fill.solid()
            box.fill.fore_color.rgb = RGBColor(180, 200, 255)
            box.line.fill.background()

            tf = box.text_frame
            p = tf.paragraphs[0]

            p.text = step
            p.font.size = Pt(14)
            p.alignment = PP_ALIGN.CENTER

            shapes.append(box)

        # Connect steps
        for i in range(1, len(shapes)):
            slide.shapes.add_connector(
                MSO_CONNECTOR.STRAIGHT,
                shapes[i - 1].left + shapes[i - 1].width,
                shapes[i - 1].top + shapes[i - 1].height / 2,
                shapes[i].left,
                shapes[i].top + shapes[i].height / 2,
            )

    def _download_image(self, query):
        """
        Download an image from Unsplash based on query.

        Args:
            query (str): Search keyword

        Returns:
            str or None: Path to downloaded image or None if failed
        """
        try:
            url = f"https://source.unsplash.com/600x400/?{query}"
            res = requests.get(url, timeout=5)

            if res.headers.get("content-type", "").startswith("image/"):
                path = f"temp_{query.replace(' ', '_')}.jpg"
                with open(path, "wb") as f:
                    f.write(res.content)
                return path
        except:
            return None

    def save(self):
        """
        Save the presentation to the output directory.

        Returns:
            str: File path of the saved presentation
        """
        output_dir = "output"

        if not os.path.exists(output_dir):
            os.makedirs(output_dir)

        timestamp = datetime.now().strftime("%Y%m%d_%H%M%S")
        path = os.path.join(output_dir, f"output_{timestamp}.pptx")

        self.prs.save(path)

        # Clean temporary images
        for f in glob.glob("temp_*.jpg"):
            try:
                os.remove(f)
            except:
                pass

        return path